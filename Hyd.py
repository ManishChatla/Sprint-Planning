from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Slide layout
title_layout = prs.slide_layouts[0]
bullet_layout = prs.slide_layouts[1]

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_detailed_slide(title, content):
    slide = prs.slides.add_slide(bullet_layout)
    slide.shapes.title.text = title
    textbox = slide.placeholders[1]
    textbox.text = content

# Add slides with detailed information
add_title_slide("Hyderabad – The City of Pearls and Progress", "A Deep Dive into My Hometown")

add_detailed_slide("A Glimpse into History", 
    "Hyderabad was founded in 1591 by Muhammad Quli Qutb Shah, the fifth ruler of the Qutb Shahi dynasty. "
    "It served as a cultural and political center for centuries, later becoming the capital of the Nizams—one of the richest princely states under British rule. "
    "The Nizams were globally known for their immense wealth, palaces, and patronage of the arts."
)

add_detailed_slide("Cultural Heritage", 
    "Hyderabad reflects a fascinating mix of Islamic, Persian, Mughal, and Telugu cultures. "
    "Its architectural marvels include the iconic Charminar, Golconda Fort, and the Chowmahalla Palace. "
    "It’s also home to the Qutb Shahi Tombs and beautiful mosques, showing its deep-rooted Islamic heritage. "
    "Classical dance, traditional art, and Urdu poetry continue to thrive."
)

add_detailed_slide("Language & Lifestyle", 
    "Hyderabad is a multilingual city where Telugu, Urdu, Hindi, and English are widely spoken. "
    "The famous 'Hyderabadi Urdu' accent is instantly recognizable. "
    "The lifestyle is a unique blend of traditional values and a modern, tech-driven urban culture, with people known for their warmth and hospitality."
)

add_detailed_slide("Famous Cuisine", 
    "The city is world-renowned for its Hyderabadi Biryani—a rich, flavorful rice dish cooked with spices and meat. "
    "Other iconic dishes include Haleem (a wheat-meat dish eaten during Ramadan), Double Ka Meetha (bread pudding), and Mirchi Ka Salan. "
    "Irani chai and Osmania biscuits are a staple in old-style cafés, especially around Charminar."
)

add_detailed_slide("What Hyderabad is Famous For", 
    "Known as the City of Pearls due to its historic pearl and diamond trade, Hyderabad is also a global IT hub with a major presence of companies like Microsoft, Google, and Amazon. "
    "It is home to Tollywood (Telugu film industry) and Ramoji Film City—the world’s largest integrated film studio. "
    "It also has a strong pharmaceutical and biotechnology industry."
)

add_detailed_slide("Modern Infrastructure", 
    "Hyderabad boasts excellent infrastructure: a well-connected Outer Ring Road (ORR), efficient Metro Rail, and the Rajiv Gandhi International Airport, one of India’s top-rated airports. "
    "HiTec City and Gachibowli are major IT and financial districts with cutting-edge facilities."
)

add_detailed_slide("Population & Demographics", 
    "With a population of around 10 million, Hyderabad is India's 6th largest city. "
    "It is religiously and culturally diverse, with a harmonious mix of Hindus, Muslims, Christians, and others. "
    "The city has a growing population of educated youth and professionals, making it dynamic and progressive."
)

add_detailed_slide("Festivals & Celebrations", 
    "Hyderabad celebrates both Hindu and Muslim festivals with grandeur. "
    "Ganesh Chaturthi, Diwali, Eid, and Bonalu are widely celebrated, often involving entire neighborhoods. "
    "The city becomes vibrant and colorful, showing true unity in diversity."
)

add_detailed_slide("Notable Places to Visit", 
    "Top attractions include the Charminar, Golconda Fort, Hussain Sagar Lake with its giant Buddha statue, Birla Mandir, and the Salar Jung Museum. "
    "Laad Bazaar is famous for bangles and pearls, while Shilparamam showcases local crafts and cultural performances."
)

add_detailed_slide("Climate & Area", 
    "Hyderabad has a tropical wet and dry climate with hot summers (March–June), a monsoon season (July–September), and mild winters (November–February). "
    "Average summer temperatures reach 35–40°C, while winters are pleasant around 20°C. "
    "The city's area is ~650 sq km, and Greater Hyderabad Metropolitan Region spans over 1200+ sq km."
)

add_detailed_slide("Education & Institutions", 
    "Hyderabad is a center of excellence in education. Key institutions include IIT Hyderabad, IIIT Hyderabad, University of Hyderabad, and Osmania University. "
    "It also houses NALSAR (law), BITS Pilani Hyderabad campus, and several reputed international schools. "
    "It is also known for competitive exam coaching centers (e.g., for UPSC, IIT-JEE)."
)

add_detailed_slide("Cost of Living", 
    "Hyderabad is one of the most affordable metro cities in India. "
    "Housing, transportation, and food are reasonably priced. "
    "The city offers a comfortable lifestyle whether you're a student, a working professional, or a family. "
    "There are both budget-friendly and luxury living options available."
)

add_detailed_slide("Fun Facts", 
    "• Declared a UNESCO Creative City of Gastronomy in 2019\n"
    "• Ramoji Film City can accommodate 50+ film units at once\n"
    "• The Nizam's jewels include the famous Jacob Diamond\n"
    "• Hyderabadi slang and chai culture are iconic and beloved"
)

add_detailed_slide("Why Visit Hyderabad?", 
    "Hyderabad offers a rare combination of historic elegance and modern development. "
    "Visitors can enjoy cultural richness, architectural wonders, diverse cuisine, and tech-city vibes—all in one place. "
    "It’s a city that welcomes everyone with a big heart... and an even bigger biryani!"
)

add_detailed_slide("Thank You!", 
    "Thank you for exploring Hyderabad with me. "
    "Hope this presentation gave you insights into the history, charm, and modern vibe of my hometown. "
    "Looking forward to hosting you someday in the City of Pearls!"
)

# Save presentation
detailed_path = "/mnt/data/Hyderabad_Presentation_Detailed.pptx"
prs.save(detailed_path)
detailed_path
